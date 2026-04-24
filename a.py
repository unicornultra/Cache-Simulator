from pptx import Presentation

# Create a PowerPoint presentation
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Security Analysis of UPI Protocols Using the Tamarin Prover"
slide.placeholders[1].text = "Final Year Project (BTP)\nSahil Kumar & Yash Joshi\nDepartment of CSE\nNovember 29, 2025"

# Slide 2 - Motivation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Motivation"
slide.placeholders[1].text = (
    "• UPI enables nation-wide real-time fund transfer\n"
    "• Highly critical due to financial sensitivity\n"
    "• Need to ensure strong security against active attackers\n"
    "• Formal verification helps detect flaws missed by manual review"
)

# Slide 3 - Scope & Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Scope & Objectives"
slide.placeholders[1].text = (
    "• Model key UPI flows using Tamarin\n"
    "• Verify secrecy & authentication\n"
    "• Check replay-resistance & MITM resilience\n"
    "• Identify vulnerabilities & propose mitigations"
)

# Slide 4 - UPI Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "UPI Overview"
slide.placeholders[1].text = (
    "• Actors: Payer, Payee, PSP/Bank, NPCI\n"
    "• Real-time collect/pay flows\n"
    "• Uses encryption, MACs, nonces & session keys"
)

# Slide 5 - Formal Verification & Tamarin
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Formal Verification & Tamarin Prover"
slide.placeholders[1].text = (
    "• Symbolic protocol verification\n"
    "• Multiset-rewriting rules for modeling\n"
    "• Automated attack discovery & property proofs"
)

# Slide 6 - Methodology: Modeling
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Modeling Approach"
slide.placeholders[1].text = (
    "• Focus on essential protocol operations\n"
    "• Encode messages as rewrite rules\n"
    "• Dolev-Yao adversary with full network control\n"
    "• Specify secrecy, authentication & replay safety"
)

# Slide 7 - Methodology: Workflow
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Tamarin Workflow"
slide.placeholders[1].text = (
    "• Write protocol theory (.spthy)\n"
    "• Encode lemmas for properties\n"
    "• Run Tamarin CLI/GUI\n"
    "• Analyze attacks or proofs"
)

# Slide 8 - Progress So Far
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Progress So Far"
slide.placeholders[1].text = (
    "• Completed Tamarin learning phase\n"
    "• Built baseline UPI collect/pay models\n"
    "• Initial secrecy & authentication lemmas added\n"
    "• Tool found potential weaknesses — analyzing"
)

# Slide 9 - Preliminary Observations
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Preliminary Observations"
slide.placeholders[1].text = (
    "• Missing nonces → replay attack possibilities\n"
    "• Binding identities improves security proofs\n"
    "• Model detail strongly affects results"
)

# Slide 10 - Limitations
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Limitations"
slide.placeholders[1].text = (
    "• Some attacks may be due to abstraction\n"
    "• Complex proofs require manual guidance"
)

# Slide 11 - Future Work
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Future Work"
slide.placeholders[1].text = (
    "• Model additional UPI flows including NPCI\n"
    "• Reduce spurious attacks via tighter rules\n"
    "• Document valid attacks & mitigations\n"
    "• Final evaluation & results presentation"
)

# Slide 12 - Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
slide.placeholders[1].text = (
    "Formal verification boosts confidence in UPI\n"
    "Tamarin helps uncover protocol vulnerabilities\n"
    "Next: complete verification & security analysis"
)

# Save presentation
file_path = "/mnt/data/UPI_Tamarin_Presentation.pptx"
prs.save(file_path)

file_path
